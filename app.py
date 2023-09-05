import streamlit as st
from yahooquery import Ticker
import pandas as pd
import plotly.express as px
from pptx import Presentation
from pptx.util import Inches
from datetime import date
from PIL import Image
import requests
import os
from io import BytesIO

# Resize Image function
def resize_image(url):
    image = Image.open(requests.get(url, stream=True).raw)
    if image.height > 140 or image.width > 220:
        container_width, container_height = 220 * 2, 140 * 2
    else:
        container_width, container_height = 220, 140
    new_image = Image.new('RGBA', (container_width, container_height))
    x, y = (container_width - image.width) // 2, (container_height - image.height) // 2
    new_image.paste(image, (x, y))
    return new_image

def add_image(slide, image, left, top, width):
    slide.shapes.add_picture(image, left=left, top=top, width=width)

def replace_text(replacements, shapes):
    for shape in shapes:
        for match, replacement in replacements.items():
            if shape.has_text_frame and (shape.text.find(match)) != -1:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    whole_text = "".join(run.text for run in paragraph.runs)
                    whole_text = whole_text.replace(str(match), str(replacement))
                    for idx, run in enumerate(paragraph.runs):
                        if idx != 0:
                            p = paragraph._p
                            p.remove(run._r)
                    if bool(paragraph.runs):
                        paragraph.runs[0].text = whole_text

def get_stock(ticker, period, interval):
    hist = ticker.history(period=period, interval=interval)
    return hist.reset_index()

def plot_graph(df, x, y, title):
    fig = px.line(df, x=x, y=y, template='simple_white', title='<b>{} {}</b>'.format(name, title))
    fig.update_layout(paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)')
    return fig

def get_financials(df, col_name, metric_name):
    metric = df.loc[:, ['asOfDate', col_name]]
    metric_df = pd.DataFrame(metric).reset_index()
    metric_df.columns = ['Symbol', 'Year', metric_name]
    return metric_df



# Streamlit setup
st.set_page_config(page_icon="🚀", page_title="Deloitte PowerPoint Generator")
path = os.path.dirname(__file__)
today = date.today()
st.title('Deloitte PowerPoint Generator!')

user_input = st.text_input(label='Enter company ticker. For example: AAPL for Apple or TSLA for Tesla')
submit = st.button(label='Generate PowerPoint slides')

# trim user input string
user_input = str(user_input.lower()).strip()

if submit and user_input == "":
    st.warning("Please enter company ticker!")

elif submit and user_input != "":
    with st.spinner('Generating  slides for you...⏳'):
        try:
            pptx = path + '//' + 'template.pptx'
            prs = Presentation(pptx)
            width = Inches(8)
            left = Inches(2.5)
            top = Inches(1)
            ticker = Ticker(user_input)
            name = ticker.price[user_input]['shortName']
            sector = ticker.summary_profile[user_input]['sector']
            industry = ticker.summary_profile[user_input]['industry']
            employees = ticker.summary_profile[user_input]['fullTimeEmployees']
            country = ticker.summary_profile[user_input]['country']
            city = ticker.summary_profile[user_input]['city']
            website = ticker.summary_profile[user_input]['website']
            summary = ticker.summary_profile[user_input]['longBusinessSummary']
            logo_url = 'https://logo.clearbit.com/' + website
            first_slide, second_slide = prs.slides[0], prs.slides[1]
            shapes_1, shapes_2 = list(first_slide.shapes), list(second_slide.shapes)

            replaces_1 = {'{company}': name, '{date}': today}
            replaces_2 = {
                '{c}': name,
                '{s}': sector,
                '{i}': industry,
                '{co}': country,
                '{ci}': city,
                '{ee}': "{:,}".format(employees),
                '{w}': website,
                '{summary}': summary
            }

            replace_text(replaces_1, shapes_1)
            replace_text(replaces_2, shapes_2)

            if requests.get(logo_url).status_code == 200:
                logo = resize_image(logo_url)
                logo.save('logo.png')
                add_image(prs.slides[1], image='logo.png', left=Inches(1.2), width=Inches(2), top=Inches(0.5))
                os.remove('logo.png')

            income_df = ticker.income_statement()
            valuation_df = ticker.valuation_measures

            stock_df = get_stock(ticker=ticker, period='5y', interval='1mo')
            stock_fig = plot_graph(df=stock_df, x='date', y='open', title='Stock Price USD')
                        # ... [previous code continues here] ...
            
            stock_fig.write_image("stock_plot.png")
            slide_3 = prs.slides[2]
            add_image(slide_3, image="stock_plot.png", left=Inches(1), width=Inches(6), top=Inches(2))
            os.remove("stock_plot.png")
            
            
            
            revenue = get_financials(income_df, "TotalRevenue", "Total Revenue")
            revenue_fig = plot_graph(df=revenue, x='Year', y='Total Revenue', title='Total Revenue USD')
            revenue_fig.write_image("revenue_plot.png")
            slide_4 = prs.slides[3]
            add_image(slide_4, image="revenue_plot.png", left=Inches(1), width=Inches(6), top=Inches(2))
            os.remove("revenue_plot.png")

            
            
            valuation = get_financials(valuation_df, "MarketCap", "Market Capitalization")
            valuation_fig = plot_graph(df=valuation, x='Year', y='Market Capitalization', title='Market Capitalization USD')
            valuation_fig.write_image("valuation_plot.png")
            slide_5 = prs.slides[4]
            add_image(slide_5, image="valuation_plot.png", left=Inches(1), width=Inches(6), top=Inches(2))
            os.remove("valuation_plot.png")

            file_path = f"{name}_slides.pptx"
            prs.save(file_path)
            with open(file_path, "rb") as f:
                 ppt_bytes = f.read()

            st.download_button(label=f"Download {name} Slides", data=ppt_bytes, file_name=f"{name}_slides.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

        
        except Exception as e:
            st.error(f"Error: {str(e)}")
            st.warning("Please make sure the ticker symbol is correct or try again later.")

